"""
Gemini OCR Converter - Convert PDF/DOC/DOCX/PPTX to Markdown using Google Gemini Vision API

Sử dụng Gemini API để OCR tài liệu (PDF scan, PDF digital, DOC, DOCX, PPTX PowerPoint)
và tạo file Markdown tương ứng.
"""

import os
import sys
import time
from pathlib import Path
from typing import Union

# Add project root to path (tools/ocr/ → tools/ → project root)
ROOT_DIR = Path(__file__).parent.parent.parent
sys.path.insert(0, str(ROOT_DIR))

from dotenv import load_dotenv
load_dotenv(ROOT_DIR / ".env")

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    print("❌ GEMINI_API_KEY không được tìm thấy trong .env")
    sys.exit(1)

import google.genai as genai

client = genai.Client(api_key=GEMINI_API_KEY)

# Thư mục tài liệu
DOCUMENTS_DIR = ROOT_DIR / "data" / "documents"
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.pptx', '.pptm'}


def count_pdf_pages(file_path: Path) -> int:
    """Đếm số trang PDF. Trả về 0 nếu không đọc được."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(file_path))
        count = len(doc)
        doc.close()
        return count
    except ImportError:
        try:
            import pypdf
            reader = pypdf.PdfReader(str(file_path))
            return len(reader.pages)
        except ImportError:
            return 0  # Không biết số trang, xử lý cả file


def split_pdf_pages(file_path: Path, start_page: int, end_page: int) -> Path:
    """Tách trang start_page đến end_page từ PDF, lưu ra file tạm."""
    import tempfile
    try:
        import fitz
        doc = fitz.open(str(file_path))
        new_doc = fitz.open()
        new_doc.insert_pdf(doc, from_page=start_page, to_page=end_page)
        temp_path = Path(tempfile.mktemp(suffix='.pdf'))
        new_doc.save(str(temp_path))
        new_doc.close()
        doc.close()
        return temp_path
    except ImportError:
        try:
            import pypdf
            reader = pypdf.PdfReader(str(file_path))
            writer = pypdf.PdfWriter()
            for i in range(start_page, min(end_page + 1, len(reader.pages))):
                writer.add_page(reader.pages[i])
            temp_path = Path(tempfile.mktemp(suffix='.pdf'))
            with open(temp_path, 'wb') as f:
                writer.write(f)
            return temp_path
        except ImportError:
            return None  # Không thể tách trang


def ocr_file_with_gemini(file_path: Path, label: str = "") -> str:
    """Upload và OCR một file (PDF chunk hoặc file gốc) bằng Gemini."""
    import shutil, tempfile
    
    temp_dir = tempfile.mkdtemp()

    # .pptm không được Gemini chấp nhận (MIME không hỗ trợ).
    # .pptm và .pptx cùng cấu trúc ZIP → đổi extension sang .pptx là đủ.
    upload_ext = '.pptx' if file_path.suffix.lower() == '.pptm' else file_path.suffix
    safe_name = f"document{upload_ext}"
    temp_file = Path(temp_dir) / safe_name
    shutil.copy2(file_path, temp_file)
    if upload_ext != file_path.suffix:
        print(f"   Auto-convert {file_path.suffix} → {upload_ext} (same ZIP structure, fix MIME)")
    
    try:
        uploaded_file = client.files.upload(file=str(temp_file))
        print(f"   Upload OK: {uploaded_file.name} {label}")
    except Exception as e:
        print(f"   Upload FAILED: {e}")
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
    
    # Đợi xử lý
    while uploaded_file.state.name == "PROCESSING":
        print(f"   Waiting for Gemini to process...")
        time.sleep(5)
        uploaded_file = client.files.get(name=uploaded_file.name)
    
    if uploaded_file.state.name == "FAILED":
        raise Exception(f"Gemini failed to process file: {uploaded_file.state.name}")
    
    # Chọn prompt phù hợp theo loại file
    ext = Path(str(uploaded_file.display_name)).suffix.lower() if hasattr(uploaded_file, 'display_name') else label
    if '.pptx' in str(label) or '.pptx' in str(ext):
        prompt = """Bạn là chuyên gia trích xuất nội dung từ file PowerPoint tiếng Việt.

NHIỆM VỤ: Trích xuất TOÀN BỘ nội dung các slide sang Markdown.

QUY TẮC:
1. Mỗi slide bắt đầu bằng: ## Slide N: [Tiêu đề slide]
2. Liệt kê đầy đủ các bullet point, nội dung văn bản trên slide.
3. Nếu slide có bảng → chuyển thành bảng Markdown.
4. Nếu slide có ghi chú (Notes) → thêm vào cuối slide dưới dạng: > 📝 Ghi chú: ...
5. Giữ nguyên số hiệu, ngày tháng, tên đơn vị.
6. KHÔNG thêm bình luận. KHÔNG bỏ sót nội dung.
7. Tiếng Việt, không dịch.

Trích xuất TOÀN BỘ nội dung ngay bây giờ:"""
    else:
        prompt = """Bạn là chuyên gia OCR văn bản pháp luật tiếng Việt.

NHIỆM VỤ: Trích xuất TOÀN BỘ nội dung tài liệu sang Markdown.

QUY TẮC:
1. Giữ nguyên cấu trúc: tiêu đề, điều khoản, khoản, điểm, bảng biểu.
2. Markdown chuẩn: # tiêu đề chính, ## Chương, ### Điều, **bold**, bảng markdown.
3. Giữ nguyên số hiệu, ngày tháng.
4. KHÔNG thêm bình luận. KHÔNG bỏ sót nội dung.
5. Tiếng Việt, không dịch.

Trích xuất TOÀN BỘ nội dung ngay bây giờ:"""

    try:
        response = client.models.generate_content(
            model="gemini-2.0-flash",
            contents=[uploaded_file, prompt],
            config={"max_output_tokens": 8192}
        )
        content = response.text
        print(f"   OCR done: {len(content)} chars {label}")
        return content
    except Exception as e:
        print(f"   OCR FAILED: {e}")
        raise
    finally:
        try:
            client.files.delete(name=uploaded_file.name)
        except:
            pass


def extract_pptx_text(file_path: Path) -> list[dict]:
    """
    Dùng python-pptx để extract text từng slide từ file .pptx/.pptm.
    Trả về list of dict: [{'slide': N, 'title': str, 'content': str, 'notes': str}]
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise ImportError("python-pptx chưa được cài. Chạy: pip install python-pptx")

    prs = Presentation(str(file_path))
    slides_data = []

    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        content_parts = []
        notes_text = ""

        # Lấy tiêu đề slide
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # Lấy nội dung từ tất cả text boxes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Bỏ qua title đã lấy ở trên
            if shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    content_parts.append(text)

        # Lấy ghi chú (Notes)
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""

        slides_data.append({
            "slide": idx,
            "title": title,
            "content": "\n".join(content_parts),
            "notes": notes_text,
        })

    return slides_data


def convert_pptx_with_gemini(file_path: Path, libreoffice_path: str = None) -> Path:
    """
    Pipeline đọc FULL nội dung slide (bảng, SmartArt, ảnh có chữ):
    1. LibreOffice: .pptm/.pptx → .pdf (giữ nguyên layout)
    2. PyMuPDF: PDF → ảnh PNG từng trang (slide)
    3. Gemini Vision: OCR từng ảnh → Markdown
    """
    import subprocess, tempfile, shutil, base64

    # --- Tìm LibreOffice ---
    lo_candidates = [
        libreoffice_path,
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",
    ]
    soffice = None
    for c in lo_candidates:
        if c and (Path(c).exists() or shutil.which(c)):
            soffice = c
            break

    if not soffice:
        raise RuntimeError(
            "Không tìm thấy LibreOffice!\n"
            "Cài tại: https://www.libreoffice.org/download/\n"
            "Sau đó chạy lại: python tools/ocr/gemini_ocr.py --files <tên_file>"
        )

    temp_dir = Path(tempfile.mkdtemp())
    try:
        # Bước 1: PPTM/PPTX → PDF
        print(f"   [1/3] LibreOffice convert {file_path.suffix} → PDF ...")
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", str(temp_dir), str(file_path)],
            capture_output=True, text=True, timeout=300
        )
        pdf_files = list(temp_dir.glob("*.pdf"))
        if not pdf_files:
            raise RuntimeError(f"LibreOffice không tạo được PDF.\nSTDERR: {result.stderr[:500]}")
        pdf_path = pdf_files[0]
        print(f"   PDF OK: {pdf_path.name}")

        # Bước 2: PDF → ảnh PNG từng slide
        print(f"   [2/3] Render từng slide thành ảnh PNG (150 DPI)...")
        try:
            import fitz
        except ImportError:
            raise ImportError("PyMuPDF chưa cài. Chạy: pip install pymupdf")

        doc = fitz.open(str(pdf_path))
        total_slides = len(doc)
        print(f"   Tổng slides: {total_slides}")

        slide_images = []
        for page_num in range(total_slides):
            page = doc[page_num]
            mat = fitz.Matrix(150/72, 150/72)
            pix = page.get_pixmap(matrix=mat)
            slide_images.append((page_num + 1, pix.tobytes("png")))
        doc.close()

        # Bước 3: Gemini Vision OCR từng batch ảnh
        print(f"   [3/3] Gemini Vision OCR...")
        all_markdown = []
        BATCH = 5

        for i in range(0, total_slides, BATCH):
            batch_imgs = slide_images[i:i+BATCH]
            batch_end = i + len(batch_imgs)
            label = f"[slides {i+1}-{batch_end}/{total_slides}]"
            print(f"   OCR {label}...")

            contents = []
            for slide_num, img_bytes in batch_imgs:
                contents.append({
                    "inline_data": {
                        "mime_type": "image/png",
                        "data": base64.b64encode(img_bytes).decode()
                    }
                })
                contents.append(f"--- Slide {slide_num} ---")

            prompt = (
                f"Bạn là chuyên gia OCR slide PowerPoint tiếng Việt.\n\n"
                f"Dưới đây là {len(batch_imgs)} ảnh slide ({i+1}→{batch_end}/{total_slides}).\n"
                "Với TỪNG slide:\n"
                "1. Bắt đầu bằng: ## Slide N: [Tiêu đề]\n"
                "2. Trích xuất TOÀN BỘ: text, bảng (→ Markdown table), danh sách, số liệu.\n"
                "3. SmartArt/sơ đồ → mô tả dạng danh sách hoặc bảng.\n"
                "4. Giữ nguyên ngôn ngữ, số liệu, tên người, SĐT.\n"
                "5. Slide trắng hoặc chỉ ảnh trang trí → ghi: (Slide hình ảnh/trang trí)\n"
                "6. KHÔNG thêm lời dẫn, lời kết.\n\nOutput Markdown:"
            )
            contents.append(prompt)

            try:
                response = client.models.generate_content(
                    model="gemini-2.0-flash",
                    contents=contents,
                    config={"max_output_tokens": 8192}
                )
                all_markdown.append(response.text)
                print(f"   OK {label}: {len(response.text)} chars")
            except Exception as e:
                print(f"   FAILED {label}: {e}")
                for slide_num, _ in batch_imgs:
                    all_markdown.append(f"## Slide {slide_num}\n\n(OCR thất bại)")

            if batch_end < total_slides:
                time.sleep(2)

        markdown_content = "\n\n".join(all_markdown)
        md_path = file_path.with_suffix('.md')
        md_path.write_text(markdown_content, encoding='utf-8')
        print(f"   Saved: {md_path.name} ({len(markdown_content)} chars)")
        return md_path

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def convert_with_gemini(file_path: Path, pages_per_batch: int = 10) -> Path:
    """
    Chuyển đổi PDF/DOC/DOCX sang Markdown bằng Gemini API.
    Tự động chia PDF lớn thành từng batch trang để tránh giới hạn output.
    """
    print(f"\n{'='*60}")
    print(f"Processing: {file_path.name}")
    print(f"   Size: {file_path.stat().st_size / 1024 / 1024:.1f} MB")
    print(f"{'='*60}")
    
    all_content = []

    if file_path.suffix.lower() == '.pdf':
        num_pages = count_pdf_pages(file_path)
        print(f"   Pages: {num_pages}")
        
        if num_pages > pages_per_batch:
            # Chia thành batch, OCR từng batch
            batches = range(0, num_pages, pages_per_batch)
            total_batches = len(list(batches))
            print(f"   Large PDF! Splitting into {total_batches} batches of {pages_per_batch} pages...")
            
            for batch_idx, start in enumerate(range(0, num_pages, pages_per_batch), 1):
                end = min(start + pages_per_batch - 1, num_pages - 1)
                label = f"[batch {batch_idx}/{total_batches} pages {start+1}-{end+1}]"
                print(f"\n   {label}")
                
                chunk_path = split_pdf_pages(file_path, start, end)
                if chunk_path is None:
                    print(f"   Cannot split PDF (no pypdf/fitz). Falling back to full file OCR.")
                    content = ocr_file_with_gemini(file_path, "[full file]")
                    all_content = [content]
                    break
                
                try:
                    content = ocr_file_with_gemini(chunk_path, label)
                    all_content.append(content)
                    # Thêm separator giữa các batch để dễ ghép
                    if batch_idx < total_batches:
                        all_content.append("\n\n")
                finally:
                    try:
                        chunk_path.unlink()
                    except:
                        pass
                
                if batch_idx < total_batches:
                    print(f"   Waiting 3s to avoid rate limit...")
                    time.sleep(3)
        else:
            # PDF nhỏ, xử lý cả file
            content = ocr_file_with_gemini(file_path, "[full file]")
            all_content = [content]
    else:
        # PPTX/PPTM: dùng python-pptx extract rồi Gemini format
        if file_path.suffix.lower() in ('.pptx', '.pptm'):
            return convert_pptx_with_gemini(file_path)
        # DOC/DOCX: upload file bình thường
        ext_label = f"[full file {file_path.suffix}]"
        content = ocr_file_with_gemini(file_path, ext_label)
        all_content = [content]
    
    # Ghép kết quả
    markdown_content = "".join(all_content)
    print(f"\n   Total content: {len(markdown_content)} chars")
    
    # Lưu file Markdown
    md_path = file_path.with_suffix('.md')
    md_path.write_text(markdown_content, encoding='utf-8')
    print(f"   Saved: {md_path.name}")
    
    return md_path



def scan_and_convert(force: bool = False, specific_files: list = None):
    """
    Quét thư mục data/documents/ và OCR các file chưa có .md
    
    Args:
        force: Nếu True, OCR lại cả file đã có .md
        specific_files: Danh sách tên file cụ thể cần OCR (không cần đường dẫn đầy đủ)
    """
    print("\n" + "="*60)
    print("🔎 QUÉT THƯ MỤC TÀI LIỆU")
    print(f"📂 Thư mục: {DOCUMENTS_DIR}")
    print("="*60)
    
    # Lấy danh sách file cần xử lý
    all_files = []
    if specific_files:
        # Chỉ xử lý file được chỉ định
        for name in specific_files:
            matches = list(DOCUMENTS_DIR.glob(f"*{name}*")) if '*' not in name else list(DOCUMENTS_DIR.glob(name))
            # nếu không có wildcard, tìm file khớp tên (bất kể extension)
            if not matches:
                for ext in SUPPORTED_EXTENSIONS:
                    candidate = DOCUMENTS_DIR / f"{name}{ext}"
                    if candidate.exists():
                        matches.append(candidate)
            all_files.extend([m for m in matches if m.suffix in SUPPORTED_EXTENSIONS])
        all_files = list(dict.fromkeys(all_files))  # deduplicate
    else:
        for ext in SUPPORTED_EXTENSIONS:
            all_files.extend(DOCUMENTS_DIR.glob(f"*{ext}"))
    
    # Lọc ra file chưa có .md (hoặc force re-OCR)
    needs_ocr = []
    already_done = []
    
    for f in sorted(all_files):
        md_path = f.with_suffix('.md')
        if md_path.exists() and not force:
            already_done.append(f)
        else:
            needs_ocr.append(f)
    
    # Hiển thị trạng thái
    print(f"\n📊 Tổng số file: {len(all_files)}")
    print(f"✅ Đã có .md: {len(already_done)}")
    print(f"⏳ Cần OCR: {len(needs_ocr)}")
    
    if already_done:
        print(f"\n--- Đã có .md ---")
        for f in already_done:
            print(f"  ✅ {f.name}")
    
    if not needs_ocr:
        print("\n🎉 Tất cả tài liệu đã được chuyển đổi!")
        return
    
    print(f"\n--- Cần OCR ({len(needs_ocr)} file) ---")
    for i, f in enumerate(needs_ocr, 1):
        size_mb = f.stat().st_size / 1024 / 1024
        print(f"  {i}. {f.name} ({size_mb:.1f} MB)")
    
    # Hỏi xác nhận
    print(f"\n⚠️ Sẽ OCR {len(needs_ocr)} tài liệu bằng Gemini API.")
    response = input("Bắt đầu? (y/n): ")
    if response.lower() != 'y':
        print("Đã hủy.")
        return
    
    # Xử lý từng file
    success = 0
    failed = []
    
    for i, f in enumerate(needs_ocr, 1):
        print(f"\n[{i}/{len(needs_ocr)}] Đang xử lý...")
        try:
            convert_with_gemini(f)
            success += 1
            # Nghỉ 2 giây giữa các file để tránh rate limit
            if i < len(needs_ocr):
                print("⏳ Nghỉ 2 giây tránh rate limit...")
                time.sleep(2)
        except Exception as e:
            print(f"❌ Lỗi: {e}")
            failed.append((f.name, str(e)))
    
    # Tổng kết
    print(f"\n{'='*60}")
    print(f"📊 KẾT QUẢ OCR")
    print(f"{'='*60}")
    print(f"✅ Thành công: {success}/{len(needs_ocr)}")
    if failed:
        print(f"❌ Thất bại: {len(failed)}")
        for name, err in failed:
            print(f"   - {name}: {err}")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Gemini OCR - Convert PDF/DOCX/PPTX to Markdown")
    parser.add_argument("--force", action="store_true", help="Re-OCR files that already have .md")
    parser.add_argument("--files", nargs="+", help="Specific file names (stems) to OCR, e.g. '11. So 81-2021'")
    args = parser.parse_args()
    scan_and_convert(force=args.force, specific_files=args.files)
